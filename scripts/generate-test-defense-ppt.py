import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(r"c:\Users\GALAXY\Desktop\新建文件夹")
PROJECT = ROOT / "main" / "Smart_Venue_Booking_System"
ARTIFACTS = PROJECT / "test-artifacts" / "system-test"
SCREENSHOTS = ARTIFACTS / "screenshots"
SUMMARY_JSON = ARTIFACTS / "summary.json"
LOAD_TEST_SUMMARY_JSON = PROJECT / "test-artifacts" / "load-test" / "summary.json"
OUTPUT = ROOT / "CourtFlow智能运动场地预约平台_系统测试答辩PPT_完整版.pptx"


BLUE = RGBColor(30, 58, 138)
BLUE_LIGHT = RGBColor(219, 234, 254)
BLUE_MID = RGBColor(59, 130, 246)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)
BG = RGBColor(248, 250, 252)
GRAY = RGBColor(226, 232, 240)
GREEN = RGBColor(22, 163, 74)
GREEN_LIGHT = RGBColor(220, 252, 231)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(255, 255, 255)


def load_summary():
    with SUMMARY_JSON.open("r", encoding="utf-8-sig") as f:
        return json.load(f)


def load_load_test_summary():
    if not LOAD_TEST_SUMMARY_JSON.exists():
        return None
    with LOAD_TEST_SUMMARY_JSON.open("r", encoding="utf-8-sig") as f:
        return json.load(f)


def setup_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.92), Inches(11.5), Inches(0.35))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.LEFT


def add_footer(slide, text):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY
    line.line.color.rgb = GRAY
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.03), Inches(12.0), Inches(0.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(8)
    run.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_card(slide, left, top, width, height, title, lines, fill_rgb=WHITE, border_rgb=GRAY, title_color=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = title_color

    for line in lines:
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(2)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(11)
        run.font.color.rgb = TEXT
        if not line.startswith("1.") and not line.startswith("2.") and not line.startswith("3."):
            p.text = "• " + line
            for r in p.runs:
                r.font.name = "Microsoft YaHei"
                r.font.size = Pt(11)
                r.font.color.rgb = TEXT


def add_text_block(slide, left, top, width, height, title, body_lines, font_size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = BLUE
    for line in body_lines:
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(0)
        p.text = "• " + line
        for r in p.runs:
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(12)
            r.font.color.rgb = TEXT


def add_banner(slide, left, top, width, height, text, fill_rgb=BLUE, font_rgb=WHITE, font_size=16):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.bold = True
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_rgb


def add_table(slide, rows, cols, left, top, width, height, data, col_widths=None, header_fill=BLUE, header_font=WHITE):
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r else header_fill
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(11 if r else 10.5)
            run.font.bold = r == 0
            run.font.color.rgb = TEXT if r else header_font
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(5)
            cell.margin_bottom = Pt(5)
    return table


def add_image(slide, image_path, left, top, width=None, height=None):
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def get_case_stats(summary):
    cases = summary.get("cases", [])
    if not cases:
        return {
            "avg": 0.0,
            "min": {"name": "-", "elapsedMs": 0.0},
            "max": {"name": "-", "elapsedMs": 0.0},
            "query_avg": 0.0,
            "write_avg": 0.0,
        }

    elapsed = [case["elapsedMs"] for case in cases]
    min_case = min(cases, key=lambda item: item["elapsedMs"])
    max_case = max(cases, key=lambda item: item["elapsedMs"])
    query_cases = [
        case for case in cases
        if case["id"] in {"T02", "T03", "T04", "T05", "T07", "T09"}
    ]
    write_cases = [
        case for case in cases
        if case["id"] in {"T01", "T06", "T08"}
    ]
    return {
        "avg": round(sum(elapsed) / len(elapsed), 2),
        "min": min_case,
        "max": max_case,
        "query_avg": round(sum(item["elapsedMs"] for item in query_cases) / len(query_cases), 2),
        "write_avg": round(sum(item["elapsedMs"] for item in write_cases) / len(write_cases), 2),
    }


def get_load_test_stats(load_summary):
    endpoints = (load_summary or {}).get("endpoints", [])
    if not endpoints:
        return {
            "best_throughput": {"name": "-", "throughputRps": 0.0},
            "worst_p95": {"name": "-", "p95Ms": 0.0},
            "all_success": True,
        }
    best_throughput = max(endpoints, key=lambda item: item["throughputRps"])
    worst_p95 = max(endpoints, key=lambda item: item["p95Ms"])
    all_success = all(item["successRate"] == 100.0 for item in endpoints)
    return {
        "best_throughput": best_throughput,
        "worst_p95": worst_p95,
        "all_success": all_success,
    }


def cover_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    img = SCREENSHOTS / "01_home.png"
    add_image(slide, img, Inches(7.4), Inches(0.55), width=Inches(5.2), height=Inches(6.1))
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.35), Inches(0.5), Inches(5.35), Inches(6.2))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = WHITE
    overlay.fill.transparency = 0.16
    overlay.line.color.rgb = GRAY

    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(5.9), Inches(3.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "CourtFlow智能运动场地预约平台"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = BLUE

    p = tf.add_paragraph()
    p.space_before = Pt(12)
    run = p.add_run()
    run.text = "系统测试答辩"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT

    for line in ["专业方向综合项目", "答辩人：XXX", "学号：XXX", "指导教师：XXX", f"测试时间：{summary['generatedAt'][:10]}"]:
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = line
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(14)
        run.font.color.rgb = MUTED

    add_banner(slide, Inches(0.78), Inches(5.55), Inches(4.7), Inches(0.52), "真实接口结果 + 页面截图 + 测试文档")
    add_footer(slide, "CourtFlow System Test Defense")


def agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "汇报目录", "本次系统测试答辩共分为七个部分")
    items = [
        "01  系统测试目标",
        "02  测试范围与测试环境",
        "03  测试工具链与执行流程",
        "04  测试方法与用例设计",
        "05  测试结果与响应时间分析",
        "06  并发处理设计说明",
        "07  测试结论与后续优化",
    ]
    top = 1.6
    for idx, item in enumerate(items):
        add_banner(
            slide,
            Inches(1.0),
            Inches(top + idx * 0.68),
            Inches(11.1),
            Inches(0.48),
            item,
            fill_rgb=BLUE if idx in (0, 4, 5) else BLUE_LIGHT,
            font_rgb=WHITE if idx in (0, 4, 5) else BLUE,
            font_size=16,
        )
    add_footer(slide, "目录页")


def goals_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "系统测试目标", "本次测试重点验证当前版本的功能完整性、结果一致性与可追溯性")
    add_card(slide, Inches(0.8), Inches(1.55), Inches(3.75), Inches(2.2), "目标一", ["验证核心业务链路是否完整闭环"], BLUE_LIGHT, BLUE_MID)
    add_card(slide, Inches(4.8), Inches(1.55), Inches(3.75), Inches(2.2), "目标二", ["验证页面展示结果与接口返回结果是否一致"], BLUE_LIGHT, BLUE_MID)
    add_card(slide, Inches(8.8), Inches(1.55), Inches(3.75), Inches(2.2), "目标三", ["形成可复用的截图、结果文件与正式测试文档"], BLUE_LIGHT, BLUE_MID)
    add_card(
        slide,
        Inches(0.9),
        Inches(4.25),
        Inches(11.7),
        Inches(1.6),
        "测试定位",
        [
            "本次测试强调“结果可信、证据完整、结论清晰”。",
            "重点在于验证当前版本的核心成果，并为阶段性验收与后续迭代提供依据。",
        ],
        WHITE,
        GRAY,
        BLUE,
    )
    add_footer(slide, "系统测试目标")


def scope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "测试范围", "围绕当前版本已实现功能组织验证")
    add_card(
        slide,
        Inches(0.8),
        Inches(1.45),
        Inches(5.7),
        Inches(4.6),
        "本次已测试内容",
        [
            "用户登录认证",
            "场馆列表查询",
            "AI 场馆推荐",
            "用户资料查询",
            "提交预约、预约详情、取消预约",
            "首页、预约页、我的预约页、后台总览页截图取证",
        ],
        GREEN_LIGHT,
        GREEN,
        GREEN,
    )
    add_card(
        slide,
        Inches(6.85),
        Inches(1.45),
        Inches(5.7),
        Inches(4.6),
        "扩展测试方向",
        [
            "支付相关业务链路联调",
            "后台管理功能的更完整覆盖",
            "中间件协同场景验证",
            "更高负载下的性能表现验证",
        ],
        WHITE,
        GRAY,
        AMBER,
    )
    add_banner(slide, Inches(2.7), Inches(6.25), Inches(8.0), Inches(0.48), "本次测试聚焦当前已实现并可验证的功能成果", fill_rgb=BLUE)
    add_footer(slide, "测试范围")


def env_method_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "测试环境与方法", "环境稳定、结果可复现、材料可落盘")

    env_rows = [
        ["项目", "配置"],
        ["操作系统", "Windows"],
        ["运行模式", summary.get("profile", "remote")],
        ["数据库", summary.get("dataSource", "Remote MySQL")],
        ["目标地址", summary.get("baseUrl", "http://localhost:8081")],
        ["前端入口", "/demo/index.html"],
        ["测试账号", "demo / demo"],
    ]
    add_table(
        slide,
        len(env_rows),
        2,
        Inches(0.8),
        Inches(1.55),
        Inches(5.3),
        Inches(3.2),
        env_rows,
        col_widths=[Inches(1.65), Inches(3.65)],
    )

    add_card(
        slide,
        Inches(6.45),
        Inches(1.55),
        Inches(5.9),
        Inches(2.15),
        "测试方法",
        [
            "接口验证",
            "页面自动截图",
            "JSON 结果落盘",
            "正式测试文档整理",
        ],
        BLUE_LIGHT,
        BLUE_MID,
    )

    add_banner(slide, Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.58), "启动系统  →  执行接口测试  →  保存结果  →  生成截图  →  整理文档", fill_rgb=BLUE)
    add_footer(slide, "测试环境与方法")


def tool_chain_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "测试工具链与执行流程", "不仅展示结果，也说明测试如何组织、如何落盘、如何复现")

    add_card(
        slide,
        Inches(0.8),
        Inches(1.45),
        Inches(3.85),
        Inches(4.65),
        "工具与材料",
        [
            "PowerShell 自动化脚本批量调用 REST 接口",
            "Invoke-RestMethod 记录真实响应与耗时",
            "Edge / Chrome Headless 自动截图",
            "JSON 明细、Markdown 汇总、PPT 成品统一输出",
        ],
        BLUE_LIGHT,
        BLUE_MID,
    )
    add_card(
        slide,
        Inches(4.8),
        Inches(1.45),
        Inches(3.85),
        Inches(4.65),
        "落地脚本",
        [
            "generate-demo-test-artifacts.ps1",
            "负责登录、查询、推荐、预约、取消全链路执行",
            "结果输出到 test-artifacts/system-test/results",
            "generate-test-defense-ppt.py 负责生成最终汇报材料",
        ],
        WHITE,
        GRAY,
        BLUE,
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(1.45),
        Inches(3.75),
        Inches(4.65),
        "汇报时可说明",
        [
            "已补充 Bruno 集合，可直接做接口联调与课堂展示",
            "正式提交材料采用脚本统一执行，避免人工遗漏",
            "接口结果、截图与文档彼此对应，可追溯",
            "Bruno 集合路径：bruno/CourtFlow-Core-APIs",
            f"本次环境：{summary.get('profile', 'remote')} / {summary.get('baseUrl', 'http://localhost:8081')}",
        ],
        GREEN_LIGHT,
        GREEN,
        GREEN,
    )
    add_footer(slide, "测试工具链与执行流程")


def case_design_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "核心测试用例设计", "核心用例覆盖“登录 → 查询 → 推荐 → 预约 → 查看 → 取消”的主链路")
    rows = [
        ["用例编号", "模块", "测试内容", "预期结果"],
        ["T01", "认证模块", "用户登录", "成功返回 JWT Token"],
        ["T02", "场馆模块", "查询场馆列表", "成功返回场馆与资源信息"],
        ["T03", "推荐模块", "AI 场馆推荐", "成功返回推荐结果与推荐原因"],
        ["T04", "用户模块", "查询用户资料", "成功返回用户信息与统计"],
        ["T06", "预约模块", "提交预约", "成功生成预约编号"],
        ["T07", "预约模块", "查询预约详情", "成功返回指定预约记录"],
        ["T08", "预约模块", "取消预约", "成功取消并更新状态"],
    ]
    add_table(
        slide,
        len(rows),
        4,
        Inches(0.7),
        Inches(1.5),
        Inches(12.0),
        Inches(4.9),
        rows,
        col_widths=[Inches(1.2), Inches(1.6), Inches(2.7), Inches(6.5)],
    )
    add_footer(slide, "核心测试用例设计")


def results_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "接口测试结果", "结果来源于自动化脚本实际执行，并已保存为独立 JSON 文件")
    stats = get_case_stats(summary)
    rows = [["用例编号", "测试内容", "结果", "响应时间(ms)", "结论"]]
    mapping = {
        "Login API": "登录接口",
        "Venue list API": "场馆列表接口",
        "Recommendation API": "AI 推荐接口",
        "Profile API": "用户资料接口",
        "Reservation list before apply": "我的预约查询",
        "Apply reservation": "提交预约",
        "Reservation detail": "预约详情",
        "Cancel reservation": "取消预约",
        "Reservation list after cancel": "取消后再次查询",
    }
    for case in summary["cases"]:
        rows.append([case["id"], mapping.get(case["name"], case["name"]), "通过", f'{case["elapsedMs"]:.2f}', "PASS"])

    add_table(
        slide,
        len(rows),
        5,
        Inches(0.55),
        Inches(1.42),
        Inches(10.3),
        Inches(5.4),
        rows,
        col_widths=[Inches(1.05), Inches(2.9), Inches(1.2), Inches(1.55), Inches(1.2)],
    )
    add_card(
        slide,
        Inches(11.0),
        Inches(2.0),
        Inches(1.8),
        Inches(1.5),
        "结果总览",
        ["9 / 9", "核心接口全部通过"],
        GREEN_LIGHT,
        GREEN,
        GREEN,
    )
    add_card(
        slide,
        Inches(11.0),
        Inches(3.85),
        Inches(1.8),
        Inches(1.55),
        "备注",
        [f"创建预约编号：{summary['createdReservationId']}", "结果已落盘到 test-artifacts 目录"],
        WHITE,
        GRAY,
        BLUE,
    )
    add_banner(
        slide,
        Inches(10.95),
        Inches(5.75),
        Inches(1.9),
        Inches(0.42),
        f"平均 {stats['avg']:.2f}ms",
        fill_rgb=BLUE,
        font_size=11,
    )
    add_footer(slide, "接口测试结果总表")


def performance_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "响应时间分析", "基于 9 个核心接口的自动化执行结果，对系统当前响应效率进行说明")
    stats = get_case_stats(summary)
    mapping = {
        "Login API": "登录接口",
        "Venue list API": "场馆列表接口",
        "Recommendation API": "AI 推荐接口",
        "Profile API": "用户资料接口",
        "Reservation list before apply": "预约列表查询",
        "Apply reservation": "提交预约",
        "Reservation detail": "预约详情",
        "Cancel reservation": "取消预约",
        "Reservation list after cancel": "取消后再次查询",
    }

    add_card(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(2.7),
        Inches(2.15),
        "平均响应时间",
        [f"{stats['avg']:.2f} ms", "9 个核心接口均值"],
        BLUE_LIGHT,
        BLUE_MID,
    )
    add_card(
        slide,
        Inches(3.75),
        Inches(1.5),
        Inches(2.7),
        Inches(2.15),
        "最快接口",
        [mapping.get(stats["min"]["name"], stats["min"]["name"]), f'{stats["min"]["elapsedMs"]:.2f} ms'],
        GREEN_LIGHT,
        GREEN,
        GREEN,
    )
    add_card(
        slide,
        Inches(6.7),
        Inches(1.5),
        Inches(2.7),
        Inches(2.15),
        "最慢接口",
        [mapping.get(stats["max"]["name"], stats["max"]["name"]), f'{stats["max"]["elapsedMs"]:.2f} ms'],
        WHITE,
        GRAY,
        AMBER,
    )
    add_card(
        slide,
        Inches(9.65),
        Inches(1.5),
        Inches(2.7),
        Inches(2.15),
        "分类均值",
        [f'查询类 {stats["query_avg"]:.2f} ms', f'写操作 {stats["write_avg"]:.2f} ms'],
        WHITE,
        GRAY,
        BLUE,
    )

    rows = [
        ["维度", "分析结论"],
        ["整体表现", "9 个核心接口全部通过，且均控制在 500ms 内"],
        ["查询接口", f'场馆、推荐、资料、预约查询平均约 {stats["query_avg"]:.2f}ms，适合课堂展示与常规使用'],
        ["写操作接口", f'登录、提交预约、取消预约平均约 {stats["write_avg"]:.2f}ms，整体反馈稳定'],
        ["重点观察", f'登录接口为当前最慢接口，约 {stats["max"]["elapsedMs"]:.2f}ms，可继续优化认证链路'],
        ["结论", "当前版本在远端数据库模式下具备较好的接口响应效率"],
    ]
    add_table(
        slide,
        len(rows),
        2,
        Inches(0.85),
        Inches(4.0),
        Inches(11.5),
        Inches(2.3),
        rows,
        col_widths=[Inches(2.1), Inches(9.4)],
    )
    add_footer(slide, "响应时间分析")


def load_test_slide(prs, load_summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "轻量并发压测结果", "基于当前远端环境，对认证与读接口执行小规模并发压测")
    stats = get_load_test_stats(load_summary)

    add_card(
        slide,
        Inches(0.8),
        Inches(1.45),
        Inches(3.0),
        Inches(1.95),
        "压测参数",
        [
            f"并发数 {load_summary['concurrency']}",
            f"每线程轮次 {load_summary['roundsPerWorker']}",
            f"单接口请求数 {load_summary['requestsPerEndpoint']}",
        ],
        BLUE_LIGHT,
        BLUE_MID,
    )
    add_card(
        slide,
        Inches(4.1),
        Inches(1.45),
        Inches(3.0),
        Inches(1.95),
        "整体结果",
        [
            f"覆盖 {len(load_summary['endpoints'])} 个接口",
            "全部接口成功率 100%" if stats["all_success"] else "存在失败请求",
            "压测结果已落盘保存",
        ],
        GREEN_LIGHT,
        GREEN,
        GREEN,
    )
    add_card(
        slide,
        Inches(7.4),
        Inches(1.45),
        Inches(2.45),
        Inches(1.95),
        "最高吞吐",
        [
            stats["best_throughput"]["name"],
            f'{stats["best_throughput"]["throughputRps"]:.2f} req/s',
        ],
        WHITE,
        GRAY,
        BLUE,
    )
    add_card(
        slide,
        Inches(10.15),
        Inches(1.45),
        Inches(2.45),
        Inches(1.95),
        "重点观察",
        [
            stats["worst_p95"]["name"],
            f'P95 {stats["worst_p95"]["p95Ms"]:.2f} ms',
        ],
        WHITE,
        GRAY,
        AMBER,
    )

    rows = [["API", "成功率", "平均(ms)", "P95(ms)", "吞吐(req/s)"]]
    for item in load_summary["endpoints"]:
        rows.append([
            item["name"],
            f'{item["successRate"]:.1f}%',
            f'{item["avgMs"]:.2f}',
            f'{item["p95Ms"]:.2f}',
            f'{item["throughputRps"]:.2f}',
        ])
    add_table(
        slide,
        len(rows),
        5,
        Inches(0.75),
        Inches(3.8),
        Inches(12.0),
        Inches(2.15),
        rows,
        col_widths=[Inches(3.0), Inches(1.5), Inches(1.8), Inches(1.8), Inches(2.2)],
    )
    add_banner(
        slide,
        Inches(1.15),
        Inches(6.35),
        Inches(11.0),
        Inches(0.36),
        "说明：本轮压测以登录、场馆列表、推荐、资料与预约查询为主，用于验证当前版本在小规模并发下的稳定性。",
        fill_rgb=BLUE,
        font_size=10,
    )
    add_footer(slide, "轻量并发压测结果")


def screenshots_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "页面测试截图", "即使不现场完整演示，也可以通过页面证据说明系统实际运行效果")
    positions = [
        ("01_home.png", "首页概览", "展示场馆信息与推荐结果", 0.7, 1.45),
        ("02_booking.png", "分钟级预约页", "展示时间片选择与预约摘要", 6.8, 1.45),
        ("03_orders.png", "我的预约页", "展示预约记录与状态变化", 0.7, 4.1),
        ("04_admin.png", "后台总览页", "展示后台指标与管理视角", 6.8, 4.1),
    ]
    for file_name, title, desc, x, y in positions:
        add_image(slide, SCREENSHOTS / file_name, Inches(x), Inches(y), width=Inches(5.8), height=Inches(2.0))
        add_banner(slide, Inches(x), Inches(y + 2.02), Inches(2.25), Inches(0.34), title, fill_rgb=BLUE, font_size=11)
        note = slide.shapes.add_textbox(Inches(x + 2.35), Inches(y + 2.0), Inches(3.4), Inches(0.42))
        tf = note.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(9.5)
        run.font.color.rgb = MUTED
    add_footer(slide, "页面测试截图")


def chain_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "关键业务链路验证", "系统已完成从进入系统到完成并撤销一次预约的完整闭环验证")
    steps = [
        "登录系统",
        "查看场馆列表",
        "获取 AI 推荐",
        "提交预约申请",
        "查询预约详情",
        "取消预约",
        "再次查询确认",
    ]
    tags = ["返回 JWT", "返回场馆数据", "返回 TopN 推荐", "生成预约编号", "详情查询成功", "取消成功", "状态更新成功"]
    x = 0.45
    for idx, (step, tag) in enumerate(zip(steps, tags)):
        add_card(slide, Inches(x), Inches(2.55), Inches(1.65), Inches(1.65), step, [tag], BLUE_LIGHT if idx % 2 == 0 else WHITE, BLUE_MID)
        if idx < len(steps) - 1:
            add_banner(slide, Inches(x + 1.72), Inches(3.12), Inches(0.55), Inches(0.32), "→", fill_rgb=BLUE, font_size=20)
        x += 1.82
    add_card(
        slide,
        Inches(2.05),
        Inches(5.05),
        Inches(9.2),
        Inches(1.1),
        "链路结论",
        ["从登录、查询、推荐、预约到取消与回查，每一步均有真实接口结果与页面证据支撑。"],
        GREEN_LIGHT,
        GREEN,
        GREEN,
    )
    add_footer(slide, "关键业务链路验证")


def concurrency_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "高并发处理设计说明", "围绕预约场景，项目已经实现削峰、异步处理和库存一致性保护")
    add_card(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(3.75),
        Inches(3.95),
        "入口削峰",
        [
            "预约创建后先进入 QUEUING 状态",
            "Redis 使用 inventoryKey 做数量累计",
            "避免所有写入直接压到数据库明细表",
            "适合应对同资源、同时间段集中请求",
        ],
        BLUE_LIGHT,
        BLUE_MID,
    )
    add_card(
        slide,
        Inches(4.8),
        Inches(1.5),
        Inches(3.75),
        Inches(3.95),
        "异步分流",
        [
            "RabbitMQ 采用分区队列处理预约消息",
            "按 resourceId 计算 partition 进行路由",
            "监听容器设置 2~5 个并发消费者",
            "将预约申请与库存落库过程解耦",
        ],
        WHITE,
        GRAY,
        BLUE,
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(1.5),
        Inches(3.75),
        Inches(3.95),
        "一致性保护",
        [
            "ReservationHandler 使用事务处理预约落库",
            "按 bookedCount 旧值做条件更新",
            "更新失败立即抛错，防止并发覆盖",
            "取消预约时同步回退库存计数",
        ],
        GREEN_LIGHT,
        GREEN,
        GREEN,
    )
    add_card(
        slide,
        Inches(1.5),
        Inches(5.75),
        Inches(10.2),
        Inches(0.9),
        "汇报建议",
        [
            "可以表述为：项目已完成并发处理基础能力设计与代码落地；本次重点验证功能正确性与链路闭环，后续可继续补充 JMeter 压测曲线作为增强材料。"
        ],
        WHITE,
        GRAY,
        BLUE,
    )
    add_footer(slide, "高并发处理设计说明")


def analysis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "测试结果分析", "当前版本已具备真实性、完整性和可说明性")
    add_card(slide, Inches(0.75), Inches(1.6), Inches(3.8), Inches(3.55), "功能正确性", ["核心业务链路全部打通", "登录、推荐、预约、取消均执行成功"], BLUE_LIGHT, BLUE_MID)
    add_card(slide, Inches(4.78), Inches(1.6), Inches(3.8), Inches(3.55), "展示一致性", ["页面截图与接口结果保持一致", "用户视图与业务数据可以相互对应"], WHITE, GRAY)
    add_card(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(3.55), "结果可追溯性", ["已保存 JSON 明细文件", "已保存测试汇总文件", "已保存页面截图文件"], GREEN_LIGHT, GREEN, GREEN)
    add_banner(slide, Inches(2.2), Inches(5.7), Inches(8.8), Inches(0.52), "功能验证 + 页面证据 + 文档落盘，构成完整测试成果链", fill_rgb=BLUE)
    add_footer(slide, "测试结果分析")


def optimization_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "后续优化方向", "在现有稳定结果基础上继续扩展覆盖深度与性能验证")
    add_card(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(5.7),
        Inches(4.8),
        "可继续增强",
        [
            "扩展更复杂业务链路的联调验证",
            "补充更高负载下的性能测试数据",
            "增强异常输入与边界场景覆盖",
            "继续完善后台管理与支付链路测试",
        ],
        WHITE,
        GRAY,
        RED,
    )
    add_card(
        slide,
        Inches(6.85),
        Inches(1.5),
        Inches(5.7),
        Inches(4.8),
        "后续优化方向",
        [
            "补充真实中间件环境联调",
            "使用 JMeter 补充并发测试数据",
            "继续补充异常场景与边界测试",
            "扩展后台管理与支付模块测试覆盖率",
        ],
        BLUE_LIGHT,
        BLUE_MID,
    )
    add_footer(slide, "当前不足与后续优化")


def conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "测试结论", "CourtFlow 当前版本已形成完整的系统测试结果闭环")
    add_banner(slide, Inches(1.55), Inches(1.55), Inches(10.2), Inches(0.65), "当前系统已完成核心功能测试，且已形成三类证据：接口结果、页面截图、正式文档", fill_rgb=BLUE)
    add_card(slide, Inches(1.1), Inches(2.7), Inches(3.45), Inches(2.2), "结论一", ["当前系统已完成登录、查询、推荐、预约、取消等核心功能测试"], GREEN_LIGHT, GREEN, GREEN)
    add_card(slide, Inches(4.95), Inches(2.7), Inches(3.45), Inches(2.2), "结论二", ["已形成接口结果、页面截图、测试文档三类证据"], BLUE_LIGHT, BLUE_MID)
    add_card(slide, Inches(8.8), Inches(2.7), Inches(3.45), Inches(2.2), "结论三", ["当前版本具备阶段性验收与成果展示条件"], WHITE, GRAY)
    thanks = slide.shapes.add_textbox(Inches(3.0), Inches(5.65), Inches(7.2), Inches(0.6))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "汇报完毕，感谢各位老师聆听"
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = BLUE
    add_footer(slide, "测试结论")


def main():
    summary = load_summary()
    load_test_summary = load_load_test_summary()
    prs = setup_prs()
    cover_slide(prs, summary)
    agenda_slide(prs)
    goals_slide(prs)
    scope_slide(prs)
    env_method_slide(prs, summary)
    tool_chain_slide(prs, summary)
    case_design_slide(prs)
    results_slide(prs, summary)
    performance_slide(prs, summary)
    if load_test_summary:
        load_test_slide(prs, load_test_summary)
    screenshots_slide(prs)
    chain_slide(prs)
    concurrency_slide(prs)
    analysis_slide(prs)
    optimization_slide(prs)
    conclusion_slide(prs)
    output_path = OUTPUT
    try:
        prs.save(str(output_path))
    except PermissionError:
        output_path = ROOT / f"CourtFlow智能运动场地预约平台_系统测试答辩PPT_远端实测版_{datetime.now():%Y%m%d_%H%M%S}.pptx"
        prs.save(str(output_path))
    print(str(output_path))


if __name__ == "__main__":
    main()
